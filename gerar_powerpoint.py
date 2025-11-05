#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script para gerar apresentação PowerPoint do TCC
Requer: pip install python-pptx pillow
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def criar_apresentacao():
    # Criar apresentação
    prs = Presentation()

    # Definir tamanho da slide (16:9)
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ===== SLIDE 1: CAPA =====
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco

    # Adicionar logo
    if os.path.exists('logo-univesp.png'):
        left_logo = Inches(4)
        top_logo = Inches(0.5)
        pic = slide1.shapes.add_picture('logo-univesp.png', left_logo, top_logo, width=Inches(2))

    # Título principal
    left_title = Inches(0.5)
    top_title = Inches(1.5)
    width_title = Inches(9)
    height_title = Inches(0.8)

    title_box = slide1.shapes.add_textbox(left_title, top_title, width_title, height_title)
    title_frame = title_box.text_frame
    title_frame.text = "UNIVERSIDADE VIRTUAL DO ESTADO DE SÃO PAULO\nUNIVESP"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtítulo
    subtitle_box = slide1.shapes.add_textbox(left_title, Inches(2.5), width_title, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Bacharel em Engenharia da Computação"
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Título do trabalho
    work_title_box = slide1.shapes.add_textbox(left_title, Inches(3.2), width_title, Inches(1.2))
    work_title_frame = work_title_box.text_frame
    work_title_frame.text = "Tutor Inteligente e Adaptativo em Inteligência Artificial\npara Personalização do Ensino de Matemática\nno 2º Ano do Ensino Médio"
    work_title_frame.paragraphs[0].font.size = Pt(22)
    work_title_frame.paragraphs[0].font.bold = True
    work_title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Autores
    authors_text = """Brenno Calderan Brossi, RA 221084
Camila Cristina Paiva Vieira, RA 2101658
Eder Clauber dos Santos dos Anjos, RA 1806662
Franklin Leandro Martins, RA 1401068
Gabriel Henrique Souza Ribeiro, RA 2213730
Julio Anderson Guimarães Leite, RA 2104217
Tiago Prado Moreira, RA 2109577
Vera Lucia de Oliveira Silverio, RA 2231483"""

    authors_box = slide1.shapes.add_textbox(left_title, Inches(4.5), width_title, Inches(1.5))
    authors_frame = authors_box.text_frame
    authors_frame.text = authors_text
    authors_frame.paragraphs[0].font.size = Pt(12)
    authors_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Orientador
    orientador_box = slide1.shapes.add_textbox(left_title, Inches(6.2), width_title, Inches(0.5))
    orientador_frame = orientador_box.text_frame
    orientador_frame.text = "Orientador: Prof. Alex Sandro Rodrigues Ancioto"
    orientador_frame.paragraphs[0].font.size = Pt(14)
    orientador_frame.paragraphs[0].font.bold = True
    orientador_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data
    date_box = slide1.shapes.add_textbox(left_title, Inches(6.8), width_title, Inches(0.3))
    date_frame = date_box.text_frame
    date_frame.text = "São Paulo - SP | 2025"
    date_frame.paragraphs[0].font.size = Pt(12)
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # ===== SLIDE 2: Contextualização =====
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Título e conteúdo
    title2 = slide2.shapes.title
    title2.text = "Contextualização"
    title2.text_frame.paragraphs[0].font.size = Pt(44)
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content2 = slide2.placeholders[1]
    tf2 = content2.text_frame
    tf2.text = "• Baixos índices de proficiência em Matemática no Ensino Médio (INEP, 2023)"
    p2 = tf2.add_paragraph()
    p2.text = "• Desempenho brasileiro abaixo da média internacional (PISA, 2023)"
    p2 = tf2.add_paragraph()
    p2.text = "• Metodologias tradicionais uniformizam o ensino, desconsiderando diferenças individuais"
    p2 = tf2.add_paragraph()
    p2.text = "• IA como alternativa promissora para personalização do ensino"
    p2 = tf2.add_paragraph()
    p2.text = "• Potencial para reduzir desigualdades e ampliar acesso (UNESCO, 2023)"

    for paragraph in tf2.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(12)

    # ===== SLIDE 3: Problema =====
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    title3 = slide3.shapes.title
    title3.text = "Problema de Pesquisa"
    title3.text_frame.paragraphs[0].font.size = Pt(44)
    title3.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content3 = slide3.placeholders[1]
    tf3 = content3.text_frame
    tf3.text = "Como desenvolver um tutor inteligente, baseado em IA, capaz de personalizar o ensino de Matemática para estudantes do Ensino Médio, considerando princípios éticos, pedagógicos e de acessibilidade?"

    p3 = tf3.add_paragraph()
    p3.text = ""
    p3 = tf3.add_paragraph()
    p3.text = "• Educação tradicional não atende diferenças individuais"
    p3 = tf3.add_paragraph()
    p3.text = "• Dificuldades persistentes em Matemática comprometem desempenho"
    p3 = tf3.add_paragraph()
    p3.text = "• Necessidade de abordagens flexíveis e individualizadas"
    p3 = tf3.add_paragraph()
    p3.text = "• IA pode promover experiências adaptativas e maior engajamento"

    tf3.paragraphs[0].font.size = Pt(20)
    tf3.paragraphs[0].font.bold = True
    for paragraph in tf3.paragraphs[2:]:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(10)

    # ===== SLIDE 4: Justificativa =====
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    title4 = slide4.shapes.title
    title4.text = "Justificativa"
    title4.text_frame.paragraphs[0].font.size = Pt(44)
    title4.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content4 = slide4.placeholders[1]
    tf4 = content4.text_frame
    tf4.text = "• Pessoal: Reflexões próprias sobre inovação educacional"
    p4 = tf4.add_paragraph()
    p4.text = "• Acadêmica: Contribuição com pesquisas sobre Sistemas Tutores Inteligentes"
    p4 = tf4.add_paragraph()
    p4.text = "• Profissional: Aplicável em plataformas digitais e cursos preparatórios"
    p4 = tf4.add_paragraph()
    p4.text = "• Social: Ampliar acesso a ferramentas de ensino personalizado"
    p4 = tf4.add_paragraph()
    p4.text = ""
    p4 = tf4.add_paragraph()
    p4.text = "IA como apoio ao professor, não substituição"
    p4.font.bold = True

    for paragraph in tf4.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(12)

    # ===== SLIDE 5: Objetivos =====
    slide5 = prs.slides.add_slide(prs.slide_layouts[1])
    title5 = slide5.shapes.title
    title5.text = "Objetivos"
    title5.text_frame.paragraphs[0].font.size = Pt(44)
    title5.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content5 = slide5.placeholders[1]
    tf5 = content5.text_frame
    tf5.text = "Objetivo Geral:"
    tf5.paragraphs[0].font.size = Pt(22)
    tf5.paragraphs[0].font.bold = True
    tf5.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    p5 = tf5.add_paragraph()
    p5.text = "Propor e descrever um modelo teórico de tutor inteligente e adaptativo, baseado em IA com foco na personalização do ensino de matemática no segundo ano do ensino médio."
    p5.font.size = Pt(18)

    p5 = tf5.add_paragraph()
    p5.text = ""
    p5 = tf5.add_paragraph()
    p5.text = "Objetivos Específicos:"
    p5.font.size = Pt(22)
    p5.font.bold = True
    p5.font.color.rgb = RGBColor(102, 126, 234)

    objetivos = [
        "• Mapear técnicas de IA aplicáveis à personalização",
        "• Integrar fundamentos pedagógicos (Piaget e Vygotsky)",
        "• Definir requisitos e personas de estudantes",
        "• Estruturar arquitetura conceitual do tutor adaptativo",
        "• Discutir vantagens, limitações e riscos do modelo"
    ]

    for obj in objetivos:
        p5 = tf5.add_paragraph()
        p5.text = obj
        p5.font.size = Pt(18)
        p5.space_after = Pt(8)

    # ===== SLIDE 6: Metodologia =====
    slide6 = prs.slides.add_slide(prs.slide_layouts[1])
    title6 = slide6.shapes.title
    title6.text = "Metodologia"
    title6.text_frame.paragraphs[0].font.size = Pt(44)
    title6.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content6 = slide6.placeholders[1]
    tf6 = content6.text_frame
    tf6.text = "• Tipo: Pesquisa qualitativa e exploratória"
    p6 = tf6.add_paragraph()
    p6.text = "• Revisão Bibliográfica: IEEE Xplore, SciELO, Portal CAPES"
    p6 = tf6.add_paragraph()
    p6.text = "• Fundamentos Pedagógicos: Piaget (1999) e Vygotsky (2007)"
    p6 = tf6.add_paragraph()
    p6.text = "• Experimentação: Modelos LLM executados localmente via Ollama"
    p6 = tf6.add_paragraph()
    p6.text = "• Testes: 10 estudantes com questionários de avaliação"
    p6 = tf6.add_paragraph()
    p6.text = "• Análise: Qualitativa com foco em usabilidade e adequação pedagógica"

    for paragraph in tf6.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.space_after = Pt(10)

    # ===== SLIDE 7: Arquitetura =====
    slide7 = prs.slides.add_slide(prs.slide_layouts[1])
    title7 = slide7.shapes.title
    title7.text = "Arquitetura do Tutor Inteligente"
    title7.text_frame.paragraphs[0].font.size = Pt(44)
    title7.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content7 = slide7.placeholders[1]
    tf7 = content7.text_frame
    tf7.text = "• Modelo do Estudante: Acompanha progresso e adapta conteúdo"
    p7 = tf7.add_paragraph()
    p7.text = "• Modelo Pedagógico: Baseado em Piaget e Vygotsky"
    p7 = tf7.add_paragraph()
    p7.text = "• Base de Conteúdos: Materiais do 2º ano do Ensino Médio"
    p7 = tf7.add_paragraph()
    p7.text = "• Motor de IA: LLMs locais (LLaMA 3, Mistral, Phi-3)"

    for paragraph in tf7.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(12)

    # ===== SLIDE 8: Vantagens =====
    slide8 = prs.slides.add_slide(prs.slide_layouts[1])
    title8 = slide8.shapes.title
    title8.text = "Vantagens da Execução Local (Ollama)"
    title8.text_frame.paragraphs[0].font.size = Pt(36)
    title8.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content8 = slide8.placeholders[1]
    tf8 = content8.text_frame
    tf8.text = "• Privacidade: Dados permanecem no ambiente local"
    p8 = tf8.add_paragraph()
    p8.text = "• Custo: Reduzido, sem dependência de servidores externos"
    p8 = tf8.add_paragraph()
    p8.text = "• Acessibilidade: Funciona com baixa conectividade"
    p8 = tf8.add_paragraph()
    p8.text = "• Customização: Ajustável para conteúdos específicos"
    p8 = tf8.add_paragraph()
    p8.text = "• Controle: Total sobre dados dos estudantes"

    for paragraph in tf8.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(12)

    # ===== SLIDE 9: Resultados =====
    slide9 = prs.slides.add_slide(prs.slide_layouts[1])
    title9 = slide9.shapes.title
    title9.text = "Resultados"
    title9.text_frame.paragraphs[0].font.size = Pt(44)
    title9.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content9 = slide9.placeholders[1]
    tf9 = content9.text_frame
    tf9.text = "• Média de avaliação: 4,2 a 4,8 (escala 0-5)"
    p9 = tf9.add_paragraph()
    p9.text = "• 10 usuários testados"
    p9 = tf9.add_paragraph()
    p9.text = "• Boa aceitação e desempenho satisfatório"
    p9 = tf9.add_paragraph()
    p9.text = "• Viabilidade comprovada de modelos locais"
    p9 = tf9.add_paragraph()
    p9.text = "• Potencial para reforçar aprendizado e explicações personalizadas"

    for paragraph in tf9.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(12)

    # ===== SLIDE 10: Conclusão =====
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    title10 = slide10.shapes.title
    title10.text = "Conclusão"
    title10.text_frame.paragraphs[0].font.size = Pt(44)
    title10.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content10 = slide10.placeholders[1]
    tf10 = content10.text_frame
    tf10.text = "• IA representa oportunidade valiosa para inovar práticas pedagógicas"
    p10 = tf10.add_paragraph()
    p10.text = "• Potencializa personalização, engajamento e eficácia das avaliações"
    p10 = tf10.add_paragraph()
    p10.text = "• Atua como ponte entre conhecimento formal e realidade dos estudantes"
    p10 = tf10.add_paragraph()
    p10.text = "• Necessário investir em formação docente e infraestrutura tecnológica"
    p10 = tf10.add_paragraph()
    p10.text = "• Uso ético e responsável pode tornar o ensino mais inclusivo e significativo"
    p10 = tf10.add_paragraph()
    p10.text = ""
    p10 = tf10.add_paragraph()
    p10.text = "A tecnologia, aliada ao olhar humano, tem potencial transformador"
    p10.font.bold = True
    p10.font.size = Pt(22)
    p10.font.color.rgb = RGBColor(102, 126, 234)

    for paragraph in tf10.paragraphs[:-1]:
        paragraph.font.size = Pt(20)
        paragraph.space_after = Pt(10)

    # ===== SLIDE 11: Agradecimentos =====
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    title11 = slide11.shapes.title
    title11.text = "Agradecimentos"
    title11.text_frame.paragraphs[0].font.size = Pt(44)
    title11.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)

    content11 = slide11.placeholders[1]
    tf11 = content11.text_frame
    tf11.text = "Obrigado pela atenção!"
    tf11.paragraphs[0].font.size = Pt(28)
    tf11.paragraphs[0].font.bold = True
    tf11.paragraphs[0].alignment = PP_ALIGN.CENTER

    p11 = tf11.add_paragraph()
    p11.text = ""
    p11 = tf11.add_paragraph()
    p11.text = "Banca Examinadora"
    p11.font.size = Pt(24)
    p11.font.bold = True
    p11.alignment = PP_ALIGN.CENTER

    p11 = tf11.add_paragraph()
    p11.text = ""
    p11 = tf11.add_paragraph()
    p11.text = "Estamos à disposição para perguntas"
    p11.font.size = Pt(20)
    p11.font.color.rgb = RGBColor(102, 126, 234)
    p11.alignment = PP_ALIGN.CENTER

    # Salvar apresentação
    filename = 'Apresentacao_TCC_Tutor_Inteligente.pptx'
    prs.save(filename)
    print(f"Apresentação PowerPoint criada com sucesso: {filename}")
    print(f"Total de slides: {len(prs.slides)}")

if __name__ == "__main__":
    try:
        criar_apresentacao()
    except ImportError:
        print("ERRO: Biblioteca python-pptx não instalada.")
        print("Execute: pip install python-pptx")
    except Exception as e:
        print(f"ERRO ao criar apresentação: {e}")
        import traceback
        traceback.print_exc()
